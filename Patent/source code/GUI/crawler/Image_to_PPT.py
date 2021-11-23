'''
version: September 23, 2021 15:05 PM
Last revision: October 22, 2021 09:00 AM

Author : ITING.TU / Chao-Hsuan Ko
'''

import os
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from PIL import Image
#from pathlib import Path
from os.path import join


class Image_to_PPT:

    def __init__(self, getId):
        self.getId = getId
        # main function
        self.PPT(self.getId)

    # 新增專利資訊到PPT
    def add_information(self, fpath, blank_slide, prs):
        slide = prs.slides.add_slide(blank_slide)
        # 設定新增文字框的位置(letf, top)以及大小(width, height)
        left, top, width, height = Cm(0), Cm(0), Cm(12), Cm(3.6)
        # 新增文欄位落
        new_paragraph = slide.shapes.add_textbox(left=left, top=top, width=width, height=height).text_frame
        # 讀取檔案
        lines = list()
        with open(fpath, 'r', encoding='utf-8') as f:
            lines = f.read().splitlines()
        # 設定文字大小
        new_paragraph.paragraphs[0].font.size = Pt(15)
        # 設定段落內容
        info = ""
        for line in lines:
            info += line +'\n'
        new_paragraph.paragraphs[0].text = info
        return 0

    def perprocess_claim(self, lines):
        text = lines.split(",")
        claim = ""
        for sentence in text:
            claim += sentence + "\n"
        return claim

    def add_keyword(self, content, slide, getId):
        # 設定新增文字框的位置(letf, top)以及大小(width, height)
        left, top, width, height = Cm(0), Cm(0), Cm(12), Cm(3.6)
        # 新增文欄位落
        new_paragraph = slide.shapes.add_textbox(left=left, top=top, width=width, height=height).text_frame
        # 設定段落內容
        text = ''
        # ['300 display panel', '310 display circuit board', '320 display driving unit', '500 vibration generating device', '700 middle frame', '800 main circuit board', '820 camera device', '900 lower cover']
        content = content.replace("[", "")
        content = content.replace("]", "")
        content = content.replace("'", "")
        allText = content.split(',')

        # mapping {id}_idtext_Mapping / {id}_keyword_postprocessing
        ROOT_PATH = './output/'
        mapping_path = join(ROOT_PATH + str(getId) + '/' + str(getId) + '_keyword_postprocessing.txt')
        fo = open(mapping_path)
        first_line = fo.readlines()
        idx = []
        content = []
        for allcontent in first_line:
            aa = allcontent.split(',')
            idx.append(aa[0])
            content.append(aa[1])
        fo.close()
        for con in allText:
            kk = con.replace(" ", "")
            if kk in idx:
                #text += con.strip() + '\n'
                if not kk == 'FIG':
                    text += kk + ' ' + content[idx.index(kk)]
        new_paragraph.paragraphs[0].text = text
        # 設定文字大小
        new_paragraph.paragraphs[0].font.size = Pt(15)
        # 設定文字顏色
        new_paragraph.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)    # #000080
        #return 0

    def add_claim(self, fpath, blank_slide, prs):
        # 設定新增文字框的位置(letf, top)以及大小(width, height)
        left, top, width, height = Cm(0), Cm(0), Cm(25.5), Cm(3.6)
        lines = list()
        # 讀取檔案
        with open(fpath, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        main_claim_list = []
        dept_claim_list = []
        #判斷主claim和附屬claim
        for i in range(len(lines)):
            dept_claim = ""
            if lines[i]!= "":
                if(i % 2 == 0):
                    main_claim_list.append(lines[i])
                else:
                    claim = self.perprocess_claim(lines[i])
                    dept_claim = claim
                    dept_claim_list.append(dept_claim)
        for i in range(len(main_claim_list)):
            slide = prs.slides.add_slide(blank_slide)
            # main claim
            textbox = slide.shapes.add_textbox(left=left, top=top, width=width, height=height).text_frame
            # 文本框的文字自動對齊
            textbox.word_wrap = True
            textbox.paragraphs[0].text = main_claim_list[i]
            # font color
            textbox.paragraphs[0].font.color.rgb = RGBColor(16, 37, 63)  
            # dependent claim
            textbox2 = slide.shapes.add_textbox(left=left, top=Cm(10), width=width, height=height).text_frame
            textbox2.word_wrap = True
            textbox2.paragraphs[0].text = dept_claim_list[i]
            # font color
            textbox2.paragraphs[0].font.color.rgb = RGBColor(55, 96, 146)  
        return 0

    def add_pictures(self, image_folder_path, image_file, slide):
        # find the image
        image_path = image_folder_path + image_file
        img = Image.open(image_path)
        # modify image size(等比例)
        if image_file == "/pdf.jpg":
            img.thumbnail((640, 640))
        #modify image size
        else:
            img = img.resize((700,500))
        img.save(image_path)
        # add image to ppt
        slide.shapes.add_picture(image_file=image_path, left=1, top=1)
        #return 0

    def PPT(self, getId):
        picfile = []
        textmapping = []

        # get mapping list
        mapping_Path = './output/' + str(getId) +'/' + str(getId) + '_idtext_Mapping.txt'
        idtext_file = open(mapping_Path, 'r', encoding="utf-8")
        for line in idtext_file:
            temp = line.split("@@")
            picfile.append(temp[0])
            textmapping.append(temp[1])

        prs = Presentation()

        image_folder_path = './output/' + str(getId) + '/images_temp/'
        image_file_list = os.listdir(image_folder_path)
        #image檔案排序
        image_file_list.sort(key=lambda x:int(x[:-4]))
        save_path = './output/' + getId + '/' + getId + '.pptx'
        info_path = './output/' + getId + '/information.txt'
        claim_path = './output/' + getId + '/claim.txt'
        pdf_img_path = './output/' + getId
        
        fileNum = []
        for file in image_file_list:
            basename = os.path.basename(file)
            filename = os.path.splitext(basename)[0]
            fileNum.append(filename)

        # SORT
        list1 = [int(x) for x in fileNum]
        list1.sort()
        # a new blank page in ppt
        blank_slide = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide)
        self.add_pictures(pdf_img_path,'/pdf.jpg' ,slide)
        self.add_information(info_path, blank_slide, prs)
        self.add_claim(claim_path, blank_slide, prs)
        # 所有圖片加入PPT
        for image_file in image_file_list:
            # add a new blank page to ppt
            slide = prs.slides.add_slide(blank_slide)
            indexNum = picfile.index(image_file)
            # 加入圖片
            self.add_pictures(image_folder_path, image_file, slide)
            # 加入圖片中關鍵字
            self.add_keyword(textmapping[indexNum], slide, getId)

        # save pptx file
        prs.save(save_path)
